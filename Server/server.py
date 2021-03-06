#!/usr/bin/env python
# pylint: disable=invalid-name

"""
CherryPy-based webservice is used for providing services to finding the best candidate.
The server provides a set of APIs to the application side to consume the services online and locally.
Summarisation feature from Pyteaser to reduce the amount of unnecessary readings for employers, 
Keyword extraction feature from textrazor to provide a set of keywords with the corresponding category,
and combination of summarisation and keyword extraction to provide a summary that is based on the given keyword(s).
"""

"""
List of libraries imported
"""
from cStringIO import StringIO
import codecs
import json
import os, os.path
import random
import string
import threading
import subprocess
import cherrypy
from cherrypy.process import plugins
import docx
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.pdfpage import PDFPage
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from cStringIO import StringIO
from pptx import Presentation
import pyteaser
import textrazor
import os
import re
from unidecode import unidecode

# import pyteaser
def worker():
    """Background Timer that runs the hello() function every 5 seconds
    TODO: this needs to be /optimized. I don't like creating the thread
    repeatedly.fixed
    """

    while True:
        t = threading.Timer(10.0, hello)
        t.start()
        t.join()


def hello():
    """Output 'Server running' on the console"""

    print ('Server running...')

class MyBackgroundThread(plugins.SimplePlugin):
    """CherryPy plugin to create a background worker thread"""

    def __init__(self, bus):
        plugins.SimplePlugin.__init__(self, bus)

        self.t = None

    def start(self):
        """Plugin entrypoint"""

        self.t = threading.Thread(target=worker)
        self.t.daemon = True
        self.t.start()

    # Start at a higher priority that "Daemonize" (which we're not using
    # yet but may in the future)
    start.priority = 85

"""
List of APIs exposed to client side.
"""
class APIController(object): \
        # pylint: disable=too-few-public-methods

    """Controller for fictional "nodes" webservice APIs"""

    """First page rendered on the client side"""
    # #     @cherrypy.tools.json_out()
    def upload(self):
        # Regular request for '/nodes' URI
        return file('./Public/html/index.html')

    """Test API for summary"""
    @cherrypy.expose
    def test(self):
        return file("./Public/html/summaries.html")

    """Page for displaying keywords for the files uploaded. Each keyword shown on the UI can be clicked.
    Clicking the keyword transits to KeyWordSearch page to show the summaries based on the clicked keyword"""
    @cherrypy.expose
    def extractPage(self):
        return file("./Public/html/ExtractText.html")

    """Page for providing features such as uploading files and summarising texts. 
    Upload button allows users to upload multiple files at once to retrieve texts in the files.
    Summarise button summarises into 3 sentences using the texts retrieved from the files"""
    @cherrypy.expose
    def newcv(self):
        return file("./Public/html/select-file.html")

    """Page for displaying keyword and normal summaries.
    The page provides a text box for users to type in their own keywords to summarise based on the keywords"""
    @cherrypy.expose
    def keywordsearch(self):
        return file("./Public/html/KeyWordSearch.html")

    """API for extracting texts from files including doc, docx, ppt, pptx, pdf"""
    def result(self, myFile):
        out = """<html>
        <body>
            <h1 style="text-align:center;">Summary</h1>
            <div style="border:1px solid #8a8a8a;border-radius: 5px;padding: 10px; max-width:600px; margin:0 auto;">%s</div>
        </body>
        </html>"""
        
        upload_file = myFile.filename
        
        #Get file instance
        with open(upload_file, 'wb') as out:
            while True:
                data = myFile.file.read(8192)
                if not data:
                    break
                out.write(data)

        #Extract texts
        if os.path.splitext(myFile.filename)[1] == '.pdf':
            data = convertPdf(upload_file)
        elif os.path.splitext(myFile.filename)[1] == '.pptx':
            data = convertPptx(upload_file)
        elif os.path.splitext(myFile.filename)[1] == '.docx':
            data = convertDocx(upload_file)
        elif os.path.splitext(myFile.filename)[1] == '.doc':
            data = convertDocxx(upload_file)
            os.remove(os.path.splitext(myFile.filename)[0] + '.docx')
        elif os.path.splitext(myFile.filename)[1] == '.ppt':
            data = convertPptxx(upload_file)
            os.remove(os.path.splitext(myFile.filename)[0] + '.pptx')
        else:
            data = "Invalid file type!"

        # Remove this line if you do not want to remove files from temp_files
        os.remove(myFile.filename)

        return data

    """API to summarise text based on default keywords or based on the given keyword(s) (given from the user)"""
    def fetchFilteredSummaries(self, data, keywords):
        return pyteaser.Summarize(data, keywords)
    
    """API to extract keywords"""
    def extractKeywords(self, data):
        return pyteaser.extract_keywords(data)

"""Function to extract texts from docx file"""
def convertDocx(path):
    document = docx.Document(path)
    docText = ''.join([
        (unidecode(unicode(paragraph.text))) for paragraph in document.paragraphs])
    return docText

"""Function to extract texts from doc file"""
def convertDocxx(path):
    for filename in  os.listdir(os.getcwd()):
        if filename.endswith('.doc'):
            subprocess.call(['soffice', '--headless', '--convert-to', 'docx', filename])
    document = docx.Document(path[:-4] + ".docx")
    docText = ''.join([
        (unidecode(unicode(paragraph.text))) for paragraph in document.paragraphs])
    return docText

"""Function to extract texts from pptx file"""
def convertPptx(path):
    prs = Presentation(path)
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []
    full_string = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(unidecode(unicode(run.text)))

    full_string = ''.join(text_runs)
    return full_string

"""Function to extract texts from ppt file"""
def convertPptxx(path):
    for filename in  os.listdir(os.getcwd()):
        if filename.endswith('.ppt'):
            subprocess.call(['soffice', '--headless', '--convert-to', 'pptx', filename])
    prs = Presentation(path[:-4] + ".pptx")
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []
    full_string = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(unidecode(unicode(run.text)))

    full_string = ''.join(text_runs)
    return full_string

"""Function to extract texts from pdf file"""
def convertPdf(path):
    rsrcmgr = PDFResourceManager()
    retstr = StringIO()
    codec = 'utf-8'
    laparams = LAParams()
    device = TextConverter(rsrcmgr, retstr, codec=codec, laparams=laparams)
    interpreter = PDFPageInterpreter(rsrcmgr, device)

    fp = file(path, 'rb')
    for page in PDFPage.get_pages(fp):
        interpreter.process_page(page)
 
    text = retstr.getvalue()
    # Cleanup
    fp.close()
    device.close()
    retstr.close()
    text = unidecode(unicode(text, encoding = "utf-8"))
    return text

def jsonify_error(status, message, traceback, version): \
        # pylint: disable=unused-argument

    """JSONify all CherryPy error responses (created by raising the
    cherrypy.HTTPError exception)
    """

    cherrypy.response.headers['Content-Type'] = 'application/json'
    response_body = json.dumps(
        {
            'error': {
                'http_status': status,
                'message': message,
            }
        })

    cherrypy.response.status = status

    return response_body


if __name__ == '__main__':
    MyBackgroundThread(cherrypy.engine).subscribe()

    dispatcher = cherrypy.dispatch.RoutesDispatcher()

    # /nodes (GET)
    dispatcher.connect(name='upload',
                       route='/',
                       action='upload',
                       controller=APIController(),
                       conditions={'method': ['GET']})

    # /nodes (GET)
    dispatcher.connect(name='result',
                       route='/result',
                       action='result',
                       controller=APIController(),
                       conditions={'method': ['POST']})

    # /nodes (GET)
    dispatcher.connect(name='fetchFilteredSummaries',
                       route='/fetchFilteredSummaries',
                       action='fetchFilteredSummaries',
                       controller=APIController(),
                       conditions={'method': ['POST']})
    
    dispatcher.connect(name='extractKeywords',
                       route='/extractKeywords',
                       action='extractKeywords',
                       controller=APIController(),
                       conditions={'method': ['POST']})

    # /nodes (GET)
    dispatcher.connect(name='extractPage',
                       route='/extractPage',
                       action='extractPage',
                       controller=APIController(),
                       conditions={'method': ['GET']})

    # /nodes (GET)
    dispatcher.connect(name='newcv',
                       route='/newcv',
                       action='newcv',
                       controller=APIController(),
                       conditions={'method': ['GET']})

    # /nodes (GET)
    dispatcher.connect(name='keywordsearch',
                       route='/keywordsearch',
                       action='keywordsearch',
                       controller=APIController(),
                       conditions={'method': ['GET']})

    # /nodes (GET)
    dispatcher.connect(name='test',
                       route='/test',
                       action='test',
                       controller=APIController(),
                       conditions={'method': ['GET']})

    config = {
        '/': {
            'request.dispatch': dispatcher,
            'error_page.default': jsonify_error,
            'tools.sessions.on': True,
            'tools.staticdir.root': os.path.abspath(os.getcwd())
        },

        '/generator': {
            'request.dispatch': cherrypy.dispatch.MethodDispatcher(),
            'tools.response_headers.on': True,
            'tools.response_headers.headers': [('Content-Type', 'text/plain')],
        },

        '/static': {
            'tools.staticdir.on': True,
            'tools.staticdir.dir': 'Public/'
        }
    }

    cherrypy.tree.mount(root=None, config=config)
    cherrypy.server.socket_host = '0.0.0.0'
    cherrypy.server.socket_port = 80

    cherrypy.engine.start()
    cherrypy.engine.block()
